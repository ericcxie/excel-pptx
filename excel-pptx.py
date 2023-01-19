"""
Name: excel-pptx 
Created by: Eric Xie
Date: Jan. 19, 2023

"""
import win32com.client
import time
import os
import psutil
import gc
import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
import psutil


# This function is called to terminate any given program such as powerpoint or excel
def terminate_process(name):
    for process in psutil.process_iter():
        if process.name() == name:
            process.terminate()
            print(f"{name} process terminated.")
            break
    else:
        print(f"{name} process not found.")


# Returns a count of the current number of objects for a given slide starting from 0
def count_objects_on_slide(path_slides, slide_index):
    prs = Presentation(path_slides)
    slide = prs.slides[slide_index]
    count = 0
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.PLACEHOLDER:
            count += 1
    print(count)
    return count


# This function will copy and paste the table from excel into a powerpoint
def excel_pptx(path_slides, slide_index, path_excel, name_sheet):
    time.sleep(2)
    terminate_process("EXCEL.EXE")
    time.sleep(2)
    ppt_app = win32com.client.Dispatch("PowerPoint.Application")

    ppt_presentation = ppt_app.Presentations.Open(path_slides)
    ppt_slide = ppt_presentation.Slides(slide_index)

    if ppt_slide is None:
        ppt_slide = ppt_presentation.Slides.Add(1, 12)

    print("test 1 passed")

    excel_app = win32com.client.Dispatch("Excel.Application")
    time.sleep(0.5)

    print("test 2 passed")

    workbook = excel_app.Workbooks.Open(path_excel, ReadOnly=1)
    print("test 3 passed")
    time.sleep(0.5)
    worksheet = workbook.Worksheets.Item(name_sheet)

    print("test 4 passed")

    used_range = worksheet.UsedRange
    used_range.Copy()

    ppt_slide.Select()
    print("Pasting object...")
    time.sleep(2)

    ppt_app.CommandBars.ExecuteMso("PasteAsEmbedded")

    print("Saving presentation...")
    time.sleep(2)
    ppt_presentation.Save()

    print("Closing excel...")
    time.sleep(2)

    workbook.Close(SaveChanges=False)
    print("Excel closed!")
    ppt_app.Quit()
    print("Powerpoint closed!")


# This function will modify the width, height, top, left properties of the object
def resize_pptx(path_slides, slide_index, shape_index, width, height, top, left):
    prs = Presentation(path_slides)

    slide = prs.slides[slide_index]
    shape = slide.shapes[shape_index]

    shape.width = Inches(width)
    shape.height = Inches(height)
    shape.top = Inches(top)
    shape.left = Inches(left)

    prs.save("updated_presentation.pptx")
    print("Presentation saved!")


# Returns the object informations for every slide
def show_object_properties(path_slides, slide_index):
    prs = Presentation(path_slides)
    slides = prs.slides

    for slide in prs.slides:
        if slides.index(slide) + 1 == slide_index:
            print("\nslide number ", slides.index(slide) + 1)
            for shape in slide.shapes:

                # Show all objects in the slide
                print(
                    "id: %s, index: %s, height: %s, width: %s, left: %s"
                    % (
                        shape.shape_id,
                        slide.shapes.index(shape),
                        round(shape.height.inches, 2),
                        round(shape.width.inches, 2),
                        round(shape.left.inches, 2),
                    ),
                    shape.shape_type,
                    shape.name,
                )


# Specify the powerpoint path
path_slides = r"path"

# Starts from index of 1
slide_index = 5
path_excel = r"path"
name_sheet = "name"

# Change table properties
width = 9.49
height = 2.73
top = 0.98
left = 0.4

count_objects_on_slide(path_slides, 0)

# print("Starting table printing!")
excel_pptx(path_slides, slide_index, path_excel, name_sheet)

# 1 is subtracted from shape_index because it starts from index 0
shape_index = count_objects_on_slide(path_slides, 0) - 1
print(f"The pasted object has an index of: {shape_index}")
resize_pptx(path_slides, slide_index - 1, shape_index, width, height, top, left)
